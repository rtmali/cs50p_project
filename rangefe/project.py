import os
import curses
import subprocess
from datetime import datetime
import docx
import pypdf
from odf.opendocument import load
from odf.text import P
import stat
import shutil
import zipfile
import tarfile

class FileExplorer:
    def __init__(self, stdscr):
        self.stdscr = stdscr
        self.current_path = os.getcwd()
        self.file_list = []
        self.current_selection = 0
        self.selected_directory_contents = []
        self.file_info = ("", "", "", "")
        self.pop_up_active = False
        self.copied_file_path = None

    def load_files(self):
        try:
            self.file_list = os.listdir(self.current_path)
            self.file_list.sort(key=lambda x: x.lower())
        except PermissionError:
            self.file_list = []

    def load_directory_contents(self, path):
        """Load the contents of the selected directory."""
        try:
            self.selected_directory_contents = os.listdir(path)
            self.selected_directory_contents.sort(key=lambda x: x.lower())
        except PermissionError:
            self.selected_directory_contents = []

    def human_readable_size(self, size):
        for unit in ['B', 'KB', 'MB', 'GB', 'TB', 'PB']:
            if size < 1024:
                return f"{size:.2f} {unit}"
            size /= 1024

    def get_file_info(self, filename):
        full_path = os.path.join(self.current_path, filename)
        try:
            stat_info = os.stat(full_path)
            size = self.human_readable_size(stat_info.st_size)
            mtime = datetime.fromtimestamp(stat_info.st_mtime).strftime('%Y-%m-%d %H:%M:%S')
            file_type = 'Directory' if os.path.isdir(full_path) else os.path.splitext(filename)[1] or 'N/A'

            permissions = stat.filemode(stat_info.st_mode)

            return permissions, file_type, size, mtime
        except FileNotFoundError:
            return 'Unknown', 'Unknown', 'N/A', 'N/A'

    def display_file_list(self):
        self.load_files()
        max_y, max_x = self.stdscr.getmaxyx()

        left_col_width = max_x // 3
        center_col_width = max_x // 3
        right_col_width = max_x // 3

        self.stdscr.clear()

        username = os.getenv("USER") or os.getenv("USERNAME")
        hostname = subprocess.getoutput("hostname")
        header = f"{username}@{hostname}: {self.current_path}"
        header = header[:max_x - 1]  # Truncate to fit
        self.stdscr.addstr(0, 0, header, curses.color_pair(5))

        self.stdscr.addstr(1, 0, '-' * (left_col_width - 1) + '\n')
        self.stdscr.addstr(2, 0, "Parent Directories")
        self.stdscr.addstr(3, 0, '-' * (left_col_width - 1) + '\n')

        parent_dirs = [d for d in os.listdir(os.path.dirname(self.current_path)) if os.path.isdir(os.path.join(os.path.dirname(self.current_path), d))]
        for index, dirname in enumerate(parent_dirs):
            if index >= max_y - 6:
                break
            display_str = f"[+] {dirname}"
            if dirname == os.path.basename(self.current_path):
                self.stdscr.addstr(4 + index, 0, display_str, curses.A_REVERSE)
            else:
                self.stdscr.addstr(4 + index, 0, display_str)

        self.stdscr.addstr(1, left_col_width, '-' * (center_col_width - 1) + '\n')
        self.stdscr.addstr(2, left_col_width, "Current Directory Files")
        self.stdscr.addstr(3, left_col_width, '-' * (center_col_width - 1) + '\n')

        for index, filename in enumerate(self.file_list):
            if index >= max_y - 6:
                break
            full_path = os.path.join(self.current_path, filename)
            is_dir = os.path.isdir(full_path)
            prefix = '[+] ' if is_dir else '    '
            display_str = f"{prefix}{filename}"
            if index == self.current_selection:
                self.stdscr.addstr(4 + index, left_col_width, display_str, curses.A_REVERSE)
            else:
                color_pair = 1 if is_dir else 2
                self.stdscr.addstr(4 + index, left_col_width, display_str, curses.color_pair(color_pair))

        selected_file = self.file_list[self.current_selection] if self.file_list else None
        if selected_file:
            selected_full_path = os.path.join(self.current_path, selected_file)
            if os.path.isdir(selected_full_path):
                self.load_directory_contents(selected_full_path)

        self.stdscr.addstr(1, left_col_width + center_col_width, '-' * (right_col_width - 1) + '\n')
        self.stdscr.addstr(2, left_col_width + center_col_width, "Expanded Directory Contents")
        self.stdscr.addstr(3, left_col_width + center_col_width, '-' * (right_col_width - 1) + '\n')

        for index, filename in enumerate(self.selected_directory_contents):
            if index >= max_y - 6:
                break
            self.stdscr.addstr(4 + index, left_col_width + center_col_width, filename)

        self.stdscr.addstr(max_y - 4, 0, '-' * (max_x - 1))
        if selected_file:
            self.file_info = self.get_file_info(selected_file)
            info_str = f"Permissions: {self.file_info[0]} | Type: {self.file_info[1]} | Size: {self.file_info[2]} | Modified: {self.file_info[3]}"
            self.stdscr.addstr(max_y - 3, 0, info_str[:max_x - 1])

        self.stdscr.refresh()

    def preview_file(self, filepath):
        self.pop_up_active = True
        max_y, max_x = self.stdscr.getmaxyx()
        preview_win = curses.newwin(max_y - 2, max_x - 2, 1, 1)
        preview_win.border(0)

        lines = []
        try:
            if filepath.lower().endswith('.doc'):
                subprocess.run(['libreoffice', '--headless', '--convert-to', 'txt', filepath])
                with open(filepath.replace('.doc', '.txt'), 'r', encoding='utf-8') as f:
                    lines = f.readlines()
            elif filepath.lower().endswith('.docx'):
                doc = docx.Document(filepath)
                lines = [para.text for para in doc.paragraphs if para.text]

            elif filepath.lower().endswith('.odt'):
                try:
                    doc = load(filepath)

                    lines = []

                    paragraphs = doc.getElementsByType(P)
                    for paragraph in paragraphs:
                        paragraph_text = []
                        for node in paragraph.childNodes:
                            if hasattr(node, 'data'):
                                paragraph_text.append(node.data)
                        lines.append(''.join(paragraph_text))

                except Exception as e:
                    lines = [f"Error reading ODT file: {str(e)}"]

            elif filepath.lower().endswith('.pdf'):
                with open(filepath, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            lines.extend(text.splitlines())
            elif filepath.lower().endswith(('.xls', '.xlsx')):
                import pandas as pd
                df = pd.read_excel(filepath)
                lines = df.to_string(index=False).splitlines()
            elif filepath.lower().endswith(('.ppt', '.pptx')):
                from pptx import Presentation
                prs = Presentation(filepath)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                           lines.append(shape.text)
            elif filepath.lower().endswith(('.txt', '.md', '.tex', '.csv', '.py', '.java', '.c', '.cpp', '.js', '.html', '.css', '.rb', '.go', '.php', '.swift')):
                with open(filepath, 'r', encoding='utf-8') as f:
                    lines = f.readlines()
            else:
                lines = ["Unsupported file type for preview."]

        except Exception as e:
            lines = [f"Error reading file: {str(e)}"]

        for i, line in enumerate(lines[:max_y - 4]):
            preview_win.addstr(i + 1, 1, line[:max_x - 3])

        preview_win.addstr(max_y - 3, 1, "Press ESC to close preview...")

        preview_win.refresh()
        while True:
            key = preview_win.getch()
            if key == 27:
                break

        preview_win.clear()
        preview_win.refresh()
        self.pop_up_active = False

    def delete_file(self):
        """Prompt to delete the selected file or directory."""
        selected_file = self.file_list[self.current_selection]
        full_path = os.path.join(self.current_path, selected_file)
        if os.path.isdir(full_path):
            confirm = self.prompt_confirmation(f"Delete directory {selected_file} and all its contents? (y/n)")
        else:
            confirm = self.prompt_confirmation(f"Delete file {selected_file}? (y/n)")

        if confirm == 'y':
            try:
                if os.path.isdir(full_path):
                    shutil.rmtree(full_path)
                else:
                    os.remove(full_path)
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error deleting {selected_file}: {str(e)}", curses.color_pair(3))

    def prompt_confirmation(self, message):
        """Display a confirmation message and wait for user input."""
        max_y, max_x = self.stdscr.getmaxyx()
        confirm_win = curses.newwin(3, len(message) + 4, max_y // 2, (max_x - len(message) - 2) // 2)
        confirm_win.border(0)
        confirm_win.addstr(1, 1, message)
        confirm_win.refresh()

        while True:
            key = confirm_win.getch()
            if key == ord('y'):
                return 'y'
            elif key == ord('n'):
                return 'n'
            elif key == 27:
                return None

    def rename_file(self):
        """Prompt to rename the selected file or directory."""
        selected_file = self.file_list[self.current_selection]
        full_path = os.path.join(self.current_path, selected_file)
        new_name = self.prompt_input(f"Rename '{selected_file}' to: ")

        if new_name:
            new_full_path = os.path.join(self.current_path, new_name)
            try:
                os.rename(full_path, new_full_path)
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error renaming: {str(e)}", curses.color_pair(3))

    def create_directory(self):
        """Prompt to create a new directory."""
        new_dir_name = self.prompt_input("New directory name: ")
        if new_dir_name:
            new_dir_path = os.path.join(self.current_path, new_dir_name)
            try:
                os.mkdir(new_dir_path)
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error creating directory: {str(e)}", curses.color_pair(3))


    def prompt_input(self, message):
        """Display an input prompt and return the user input."""
        input_str = ""
        max_length = 20

        while True:
            max_y, max_x = self.stdscr.getmaxyx()
            window_width = max(len(message) + len(input_str) + 3, max_length)
            input_win = curses.newwin(3, window_width, max_y // 2, (max_x - window_width) // 2)
            input_win.border(0)
            input_win.addstr(1, 1, message)
            input_win.addstr(1, len(message) + 1, input_str)
            input_win.refresh()

            key = input_win.getch()
            if key in (curses.KEY_ENTER, 10, 13):
                break
            elif key == 27:
                return None
            elif key == curses.KEY_BACKSPACE or key == 127:
                input_str = input_str[:-1]
            else:
                input_str += chr(key)

        input_win.clear()
        input_win.refresh()
        return input_str.strip()

    def copy_file(self):
        """Copy the selected file or directory path."""
        selected_file = self.file_list[self.current_selection]
        self.copied_file_path = os.path.join(self.current_path, selected_file)
        self.stdscr.addstr(0, 0, f"Copied: {self.copied_file_path}", curses.color_pair(3))

    def paste_file(self):
        """Paste the copied file or directory into the current directory."""
        if self.copied_file_path:
            try:
                dest_path = os.path.join(self.current_path, os.path.basename(self.copied_file_path))
                if os.path.isdir(self.copied_file_path):
                    shutil.copytree(self.copied_file_path, dest_path)
                else:
                    shutil.copy2(self.copied_file_path, dest_path)
                self.stdscr.addstr(0, 0, f"Pasted: {dest_path}", curses.color_pair(3))
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error pasting: {str(e)}", curses.color_pair(3))

    def move_file(self):
        """Move the selected file or directory to the current directory."""
        if self.copied_file_path:
            try:
                dest_path = os.path.join(self.current_path, os.path.basename(self.copied_file_path))
                shutil.move(self.copied_file_path, dest_path)
                self.stdscr.addstr(0, 0, f"Moved: {dest_path}", curses.color_pair(3))
                self.copied_file_path = None
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error moving: {str(e)}", curses.color_pair(3))

    def move_file_with_input(self):
        """Prompt the user for a destination path to move the selected file or directory."""
        destination_path = self.prompt_input("Enter destination path: ")
        if destination_path and os.path.exists(destination_path):
            if self.copied_file_path:
                try:
                    dest_full_path = os.path.join(destination_path, os.path.basename(self.copied_file_path))
                    shutil.move(self.copied_file_path, dest_full_path)
                    self.stdscr.addstr(0, 0, f"Moved: {dest_full_path}", curses.color_pair(3))
                    self.copied_file_path = None
                except Exception as e:
                    self.stdscr.addstr(0, 0, f"Error moving: {str(e)}", curses.color_pair(3))
        else:
            self.stdscr.addstr(0, 0, "Invalid destination path.", curses.color_pair(3))

    def find(self):
        """Search for files matching a given name."""
        search_term = self.prompt_input("Enter filename to search: ")
        if search_term:
            results = [f for f in self.file_list if search_term.lower() in f.lower()]
            if results:
                self.stdscr.addstr(0, 0, "Search results:", curses.color_pair(3))
                for index, filename in enumerate(results):
                    self.stdscr.addstr(index + 1, 0, filename)
            else:
                self.stdscr.addstr(0, 0, "No files found.", curses.color_pair(3))
            self.stdscr.refresh()
            self.stdscr.getch()

    def go_to_directory(self):
        """Prompt the user for a directory path and navigate to it."""
        path = self.prompt_input("Enter directory path: ")
        if path and os.path.isdir(path):
            self.current_path = path
            self.current_selection = 0
        else:
            self.stdscr.addstr(0, 0, "Invalid directory path.", curses.color_pair(3))

    def show_file_info(self):
        """Show detailed info about the selected file."""
        selected_file = self.file_list[self.current_selection]
        full_path = os.path.join(self.current_path, selected_file)
        info = self.get_file_info(selected_file)

        info_str = f"File: {selected_file}\n"
        info_str += f"Permissions: {info[0]}\n"
        info_str += f"Type: {info[1]}\n"
        info_str += f"Size: {info[2]}\n"
        info_str += f"Modified: {info[3]}"

        max_y, max_x = self.stdscr.getmaxyx()
        info_win = curses.newwin(max_y - 2, max_x - 2, 1, 1)
        info_win.border(0)
        for i, line in enumerate(info_str.splitlines()):
            if i < max_y - 4:
                info_win.addstr(i + 1, 1, line)

        info_win.addstr(max_y - 3, 1, "Press ESC to close...")
        info_win.refresh()
        while True:
            key = info_win.getch()
            if key == 27:
                break

        info_win.clear()
        info_win.refresh()

    def compress_files(self):
        """Compress the selected file or directory into a ZIP or tar.gz file."""
        selected_file = self.file_list[self.current_selection]
        full_path = os.path.join(self.current_path, selected_file)
        output_name = self.prompt_input("Enter name for the compressed file (without extension): ")

        if output_name:
            zip_path = os.path.join(self.current_path, f"{output_name}.zip")
            try:
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    if os.path.isdir(full_path):
                        for foldername, subfolders, filenames in os.walk(full_path):
                            for filename in filenames:
                                file_path = os.path.join(foldername, filename)
                                zipf.write(file_path, os.path.relpath(file_path, full_path))
                    else:
                        zipf.write(full_path, os.path.basename(full_path))
                self.stdscr.addstr(0, 0, f"Compressed to: {zip_path}", curses.color_pair(3))
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error compressing: {str(e)}", curses.color_pair(3))

    def decompress_file(self):
        """Decompress the selected ZIP or tar.gz file."""
        selected_file = self.file_list[self.current_selection]
        full_path = os.path.join(self.current_path, selected_file)

        if selected_file.endswith('.zip'):
            try:
                with zipfile.ZipFile(full_path, 'r') as zipf:
                    zipf.extractall(self.current_path)
                self.stdscr.addstr(0, 0, f"Decompressed: {full_path}", curses.color_pair(3))
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error decompressing: {str(e)}", curses.color_pair(3))
        elif selected_file.endswith('.tar.gz'):
            try:
                with tarfile.open(full_path, 'r:gz') as tarf:
                    tarf.extractall(self.current_path)
                self.stdscr.addstr(0, 0, f"Decompressed: {full_path}", curses.color_pair(3))
            except Exception as e:
                self.stdscr.addstr(0, 0, f"Error decompressing: {str(e)}", curses.color_pair(3))
        else:
            self.stdscr.addstr(0, 0, "Selected file is not a supported archive format.", curses.color_pair(3))

    def navigate(self):
        """Main loop for navigating files."""
        while True:
            if not self.pop_up_active:
                self.display_file_list()
                key = self.stdscr.getch()
                if key in (curses.KEY_UP, ord('k')) and self.current_selection > 0:
                    self.current_selection -= 1
                elif key in (curses.KEY_DOWN, ord('j')) and self.current_selection < len(self.file_list) - 1:
                    self.current_selection += 1
                elif key in (curses.KEY_LEFT, ord('h')):
                    self.current_path = os.path.dirname(self.current_path)
                    self.current_selection = 0
                elif key in (curses.KEY_RIGHT, ord('l')):
                    if self.file_list:
                        selected_file = self.file_list[self.current_selection]
                        full_path = os.path.join(self.current_path, selected_file)
                        if os.path.isdir(full_path):
                            self.current_path = full_path
                            self.current_selection = 0

                elif key == ord('o'):
                    selected_file = self.file_list[self.current_selection]
                    self.open_file(os.path.join(self.current_path, selected_file))
                elif key == ord('p'):
                    selected_file = self.file_list[self.current_selection]
                    self.preview_file(os.path.join(self.current_path, selected_file))
                elif key == ord('d'):
                    self.delete_file()
                elif key == ord('r'):
                    self.rename_file()
                elif key == ord('n'):
                    self.create_directory()
                elif key == ord('c'):
                    self.copy_file()
                elif key == ord('v'):
                    self.paste_file()
                elif key == ord('m'):
                    self.move_file_with_input()
                elif key == ord('f'):
                    self.find()
                elif key == ord('g'):
                     self.go_to_directory()
                elif key == ord('i'):
                    self.show_file_info()
                elif key == ord('?'):
                    self.show_help()
                elif key == ord('z'):
                    self.compress_files()
                elif key == ord('e'):
                    self.decompress_file()
                elif key == ord('q'):
                    break

    def open_file(self, filepath):
        """Open the selected file with the default application."""
        try:
            subprocess.run(['xdg-open', filepath], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        except Exception as e:
            self.stdscr.addstr(0, 0, f"Error opening file: {str(e)}", curses.color_pair(3))

    def show_help(self):
        """Display help screen for keyboard shortcuts in a pop-up window."""
        self.pop_up_active = True
        help_text = [
            "Keyboard Shortcuts:",
            "---------------------",
            "Arrow Up / Arrow Down: Navigate files",
            "Arrow Left: Go to parent directory",
            "Arrow Right: Open directory",
            "o: Open file with default application",
            "p: Preview file",
            "d: Delete file/directory",
            "r: Rename file/directory",
            "m: Move file/directory",
            "n: Create new directory",
            "f: Find file/directory",
            "g: Go to a specific directory",
            "i: Show detailed information",
            "?: Show this help screen",
            "z: Compress file",
            "e: Decompress file",
            "q: Quit"
        ]

        max_y, max_x = self.stdscr.getmaxyx()
        help_win = curses.newwin(max_y - 2, max_x - 2, 1, 1)
        help_win.border(0)

        for i, line in enumerate(help_text):
            if i < max_y - 4:
                help_win.addstr(i + 1, 1, line)

        help_win.addstr(max_y - 3, 1, "Press ESC to return...")

        help_win.refresh()
        while True:
            key = help_win.getch()
            if key == 27:
                break

        help_win.clear()
        help_win.refresh()
        self.pop_up_active = False

def main(stdscr):
    curses.curs_set(0)
    curses.start_color()
    curses.init_pair(1, curses.COLOR_CYAN, curses.COLOR_BLACK)
    curses.init_pair(2, curses.COLOR_WHITE, curses.COLOR_BLACK)
    curses.init_pair(3, curses.COLOR_YELLOW, curses.COLOR_BLACK)
    curses.init_pair(5, curses.COLOR_MAGENTA, curses.COLOR_BLACK)

    explorer = FileExplorer(stdscr)
    explorer.navigate()

def load_current_directory():
    """Load the current directory files."""
    return os.listdir(os.getcwd())

def create_new_directory(name):
    """Create a new directory with the given name."""
    os.mkdir(name)

def delete_file(file_path):
    """Delete the specified file."""
    os.remove(file_path)

if __name__ == '__main__':
    curses.wrapper(main)
